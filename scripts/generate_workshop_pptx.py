"""Gera apresentacao PPTX do workshop de Observabilidade em IA."""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Paleta
# ---------------------------------------------------------------------------
PRIMARY_DARK = RGBColor(0x0D, 0x3B, 0x66)
PRIMARY_ACCENT = RGBColor(0x1A, 0x8F, 0xAD)
SECONDARY_ACCENT = RGBColor(0xE8, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2D, 0x34, 0x36)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
CODE_BG = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
TABLE_ALT = RGBColor(0xEE, 0xF5, 0xF8)

# ---------------------------------------------------------------------------
# Utilitarios
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=16,
    bold=False,
    color=DARK_TEXT,
    alignment=PP_ALIGN.LEFT,
    font_name="Calibri",
    anchor=MSO_ANCHOR.TOP,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return tf


def _add_rich_textbox(slide, left, top, width, height):
    """Add a textbox and return (textframe, para) for rich building."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    return tf


def _add_run(
    para, text, size=16, bold=False, color=DARK_TEXT, font_name="Calibri", italic=False
):
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    run.font.italic = italic
    return run


def _add_bullet(
    tf,
    text,
    level=0,
    size=16,
    color=DARK_TEXT,
    bold=False,
    space_before=Pt(4),
    space_after=Pt(2),
):
    if level == 0 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.font.bold = bold
    p.level = level
    p.space_before = space_before
    p.space_after = space_after
    return p


def _add_code_block(slide, left, top, width, height, lines, font_size=11):
    """Draw a gray rect background + monospaced code lines."""
    from pptx.util import Emu

    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = "Consolas"
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(1)
        p.space_after = Pt(1)
    return shape


def _add_rect(slide, left, top, width, height, color, line=False):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if not line:
        shape.line.fill.background()
    return shape


def _page_number(slide, num, total=24):
    _add_textbox(
        slide,
        Inches(12.2),
        Inches(7.05),
        Inches(1),
        Inches(0.35),
        f"{num}/{total}",
        font_size=9,
        color=RGBColor(0x99, 0x99, 0x99),
        alignment=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# Construcao
# ---------------------------------------------------------------------------


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    BLANK = prs.slide_layouts[6]

    def new_slide():
        return prs.slides.add_slide(BLANK)

    # =================================================================
    # SLIDE 1 — Titulo
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, PRIMARY_DARK)
    # Accent bar at top
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(1),
        Inches(2.2),
        Inches(11),
        Inches(1.2),
        "Observabilidade em IA",
        font_size=44,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
    )
    _add_textbox(
        s,
        Inches(1),
        Inches(3.5),
        Inches(11),
        Inches(0.8),
        "Workshop prático — 1 hora · MLflow GenAI",
        font_size=22,
        color=PRIMARY_ACCENT,
        alignment=PP_ALIGN.CENTER,
    )
    # thin line
    _add_rect(s, Inches(5), Inches(4.6), Inches(3.333), Inches(0.03), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(1),
        Inches(5.0),
        Inches(11),
        Inches(0.6),
        "Baseado no projeto ia-observability",
        font_size=14,
        color=RGBColor(0xAA, 0xCC, 0xDD),
        alignment=PP_ALIGN.CENTER,
    )
    _page_number(s, 1)

    # =================================================================
    # SLIDE 2 — Roteiro
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, WHITE)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.4),
        Inches(5),
        Inches(0.6),
        "Roteiro",
        font_size=32,
        bold=True,
        color=PRIMARY_DARK,
    )
    # Table
    rows, cols = 6, 3
    tbl_shape = s.shapes.add_table(
        rows, cols, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.2)
    )
    tbl = tbl_shape.table
    # Column widths
    tbl.columns[0].width = Inches(1.8)
    tbl.columns[1].width = Inches(4.0)
    tbl.columns[2].width = Inches(5.9)

    headers = ["Tempo", "Bloco", "Conteúdo"]
    data = [
        [
            "0–10 min",
            "Por que observabilidade em IA?",
            'A "caixa-preta" dos LLMs e o que muda em produção',
        ],
        [
            "10–20 min",
            "Caso de uso e problema",
            "O incidente que ninguém consegue debugar",
        ],
        ["20–25 min", "Os 4 pilares", "Tracing, Custo, Avaliação, Monitoramento"],
        ["25–50 min", "Como usar (demos ao vivo)", "Código + MLflow UI"],
        ["50–60 min", "Fechamento + Q&A", "Checklist e próximos passos"],
    ]
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_DARK

    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(15)
                p.font.color.rgb = DARK_TEXT
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.LEFT
                p.space_before = Pt(4)
                p.space_after = Pt(4)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    _page_number(s, 2)

    # =================================================================
    # SLIDE 3 — Section: Por que observabilidade em IA é importante?
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, PRIMARY_DARK)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(1),
        Inches(2.8),
        Inches(11),
        Inches(0.6),
        "1",
        font_size=18,
        color=PRIMARY_ACCENT,
        alignment=PP_ALIGN.CENTER,
    )
    _add_textbox(
        s,
        Inches(1),
        Inches(3.3),
        Inches(11),
        Inches(1),
        "Por que observabilidade em IA é importante?",
        font_size=36,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
    )
    _page_number(s, 3)

    # =================================================================
    # SLIDE 4 — A diferença entre software tradicional e IA
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, WHITE)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.4),
        Inches(11),
        Inches(0.6),
        "A diferença entre software tradicional e IA",
        font_size=28,
        bold=True,
        color=PRIMARY_DARK,
    )

    tf = _add_rich_textbox(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5))
    _add_bullet(
        tf,
        "Software tradicional: mesmo input → mesmo output",
        size=18,
        bold=True,
        space_before=Pt(8),
    )
    _add_bullet(tf, "Com LLMs isso quebra:", size=18, bold=True, space_before=Pt(16))
    for item in [
        "Não-determinismo: o mesmo prompt pode gerar respostas diferentes",
        "Caixa-preta: você não vê por que o modelo respondeu daquela forma",
        "Falha silenciosa: o modelo alucina com confiança e devolve HTTP 200",
        "Custo variável: cada chamada custa tokens que se multiplicam em escala",
        "Pipelines complexos: RAG + tools + múltiplas idas ao LLM — qual etapa falhou?",
    ]:
        _add_bullet(tf, f"• {item}", level=1, size=17, space_before=Pt(6))
    _page_number(s, 4)

    # =================================================================
    # SLIDE 5 — O que é observabilidade em IA?
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, WHITE)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.4),
        Inches(11),
        Inches(0.6),
        "O que é observabilidade em IA?",
        font_size=28,
        bold=True,
        color=PRIMARY_DARK,
    )

    # Quote box
    qbox = _add_rect(
        s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.0), SECONDARY_ACCENT
    )
    tf = qbox.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(20)
    tf.margin_right = Pt(20)
    tf.margin_top = Pt(12)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _add_run(
        p,
        "“É a capacidade de entender o estado interno de uma aplicação de IA\n"
        "a partir das suas saídas externas: traces, tokens/custo, latência e qualidade.”",
        size=16,
        italic=True,
        color=PRIMARY_DARK,
    )

    # Steps
    steps = [
        ("1", "Requisição do usuário", "O input chega na aplicação LLM"),
        ("2", "Trace completo", "Inputs, outputs, latência de cada passo"),
        ("3", "Análise", "Tokens, custo, score de qualidade"),
    ]
    y_start = Inches(2.6)
    for i, (num, title, desc) in enumerate(steps):
        y = y_start + Emu(int(Inches(1.0)) * i)
        # circle-ish number
        circ = _add_rect(s, Inches(0.8), y, Inches(0.5), Inches(0.5), PRIMARY_ACCENT)
        tf = circ.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        tf.margin_top = Pt(4)
        _add_textbox(
            s,
            Inches(1.5),
            y,
            Inches(4),
            Inches(0.5),
            title,
            font_size=18,
            bold=True,
            color=DARK_TEXT,
        )
        _add_textbox(
            s,
            Inches(1.5),
            y + Inches(0.25),
            Inches(6),
            Inches(0.4),
            desc,
            font_size=14,
            color=RGBColor(0x66, 0x66, 0x66),
        )

    # Mensagem-chave
    _add_rect(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.7), PRIMARY_DARK)
    _add_textbox(
        s,
        Inches(1.0),
        Inches(5.9),
        Inches(11.3),
        Inches(0.6),
        "Sem observabilidade, melhorar uma aplicação de IA vira chute. "
        "Com ela, vira engenharia.",
        font_size=16,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
    )

    _page_number(s, 5)

    # =================================================================
    # SLIDE 6 — Section: Caso de uso
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, PRIMARY_DARK)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(1),
        Inches(2.8),
        Inches(11),
        Inches(0.6),
        "2",
        font_size=18,
        color=PRIMARY_ACCENT,
        alignment=PP_ALIGN.CENTER,
    )
    _add_textbox(
        s,
        Inches(1),
        Inches(3.3),
        Inches(11),
        Inches(1),
        "Caso de uso e o problema que ela resolve",
        font_size=36,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
    )
    _page_number(s, 6)

    # =================================================================
    # SLIDE 7 — Cenário
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, WHITE)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.4),
        Inches(11),
        Inches(0.6),
        "O cenário",
        font_size=28,
        bold=True,
        color=PRIMARY_DARK,
    )

    # Quote box
    qbox = _add_rect(
        s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.3), SECONDARY_ACCENT
    )
    tf = qbox.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(20)
    tf.margin_right = Pt(20)
    tf.margin_top = Pt(16)
    p = tf.paragraphs[0]
    _add_run(
        p,
        "“Um cliente reclamou que o bot deu uma resposta errada às 14h32. "
        "O time não consegue reproduzir. A conta de inferência triplicou no mês. "
        "Ninguém sabe qual ferramenta está lenta nem qual versão do prompt está no ar.”",
        size=15,
        italic=True,
        color=DARK_TEXT,
    )

    _add_textbox(
        s,
        Inches(0.8),
        Inches(2.9),
        Inches(11),
        Inches(0.5),
        "Sem observabilidade, cada pergunta é impossível de responder:",
        font_size=18,
        bold=True,
        color=DARK_TEXT,
    )

    problems = [
        "“O que o usuário enviou e o que o modelo respondeu?”",
        "“Qual ferramenta o agente chamou? Ela falhou ou demorou?”",
        "“Por que a conta de tokens explodiu?”",
        "“Essa resposta estava certa? Quantas estão erradas?”",
        "“Qual versão do prompt gerou essa resposta?”",
        "“Quero auditar 100% de pagamentos, mas só 10% do resto”",
    ]
    tf = _add_rich_textbox(s, Inches(0.8), Inches(3.5), Inches(11), Inches(3.5))
    for prob in problems:
        _add_bullet(tf, f"❓ {prob}", size=15, space_before=Pt(5))

    _page_number(s, 7)

    # =================================================================
    # SLIDE 8 — O que resolve
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, WHITE)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.4),
        Inches(11),
        Inches(0.6),
        "O que cada problema resolve",
        font_size=28,
        bold=True,
        color=PRIMARY_DARK,
    )

    rows, cols = 7, 2
    tbl_shape = s.shapes.add_table(
        rows, cols, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.0)
    )
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(7.5)
    tbl.columns[1].width = Inches(4.2)

    for ci, h in enumerate(["Pergunta do time", "O que resolve"]):
        cell = tbl.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_DARK

    table_data = [
        ["O que o usuário enviou e o que o modelo respondeu?", "Tracing"],
        ["Qual ferramenta o agente chamou? Falhou ou demorou?", "Spans de tool"],
        ["Por que a conta de tokens explodiu?", "Token usage + custo"],
        ["Essa resposta estava certa? Quantas estão erradas?", "Avaliação / Judges"],
        ["Qual versão do prompt gerou isso?", "Prompt registry / versionamento"],
        ["Auditar 100% de pagamentos, mas só 10% do resto", "Sampling em produção"],
    ]
    for ri, row in enumerate(table_data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(15)
                p.font.color.rgb = DARK_TEXT
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.LEFT
                p.space_before = Pt(4)
                p.space_after = Pt(4)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            # Bold the solution column
            if ci == 1:
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True
                    p.font.color.rgb = PRIMARY_ACCENT

    _page_number(s, 8)

    # =================================================================
    # SLIDE 9 — Tool lenta no trace
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, WHITE)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.4),
        Inches(11),
        Inches(0.6),
        "Como uma tool lenta aparece num trace",
        font_size=28,
        bold=True,
        color=PRIMARY_DARK,
    )

    _add_code_block(
        s,
        Inches(0.8),
        Inches(1.3),
        Inches(11.7),
        Inches(2.2),
        [
            "@tool",
            "def check_inventory(product: str) -> str:",
            '    """Consulta o estoque disponivel de um produto."""',
            "    # Simula latencia alta seguida de falha",
            "    time.sleep(2.5)",
            "    return f\"ERRO: Timeout ao consultar estoque do produto '{product}' - API indisponivel\"",
        ],
        font_size=14,
    )

    # Highlight box
    hbox = _add_rect(
        s, Inches(0.8), Inches(3.8), Inches(11.7), Inches(1.0), SECONDARY_ACCENT
    )
    tf = hbox.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(20)
    tf.margin_right = Pt(20)
    tf.margin_top = Pt(14)
    p = tf.paragraphs[0]
    _add_run(p, "💡 ", size=18)
    _add_run(p, "No MLflow UI, esse ", size=18, color=DARK_TEXT)
    _add_run(p, "time.sleep(2.5)", size=18, bold=True, color=ACCENT_ORANGE)
    _add_run(
        p,
        " aparece como um span de 2,5s dentro do trace — "
        "o gargalo fica visível e mensurável, em vez de ser um mistério.",
        size=18,
        color=DARK_TEXT,
    )

    _page_number(s, 9)

    # =================================================================
    # SLIDE 10 — Section: Os 4 pilares
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, PRIMARY_DARK)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(1),
        Inches(2.8),
        Inches(11),
        Inches(0.6),
        "3",
        font_size=18,
        color=PRIMARY_ACCENT,
        alignment=PP_ALIGN.CENTER,
    )
    _add_textbox(
        s,
        Inches(1),
        Inches(3.3),
        Inches(11),
        Inches(1),
        "Os 4 pilares da observabilidade de LLM",
        font_size=36,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
    )
    _page_number(s, 10)

    # =================================================================
    # SLIDE 11 — 4 pilares (cards)
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, WHITE)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.4),
        Inches(11),
        Inches(0.6),
        "4 pilares da observabilidade",
        font_size=28,
        bold=True,
        color=PRIMARY_DARK,
    )

    cards = [
        (
            "🔍",
            "Tracing",
            "Registra cada passo da execução:\ninputs, outputs, latência",
            PRIMARY_DARK,
        ),
        (
            "💰",
            "Custo & Tokens",
            "Quanto cada chamada consome\ne quanto custa",
            PRIMARY_ACCENT,
        ),
        (
            "✅",
            "Avaliação",
            "Mede a qualidade das respostas\ncom judges e scorers",
            RGBColor(0x27, 0xAE, 0x60),
        ),
        (
            "🚀",
            "Produção",
            "Opera em escala: sampling,\nfeedback, sessions e users",
            ACCENT_ORANGE,
        ),
    ]

    for i, (emoji, title, desc, color) in enumerate(cards):
        x = Inches(0.5) + Inches(3.1) * i
        y = Inches(1.5)
        # Card bg
        card = _add_rect(s, x, y, Inches(2.9), Inches(4.5), WHITE, line=True)
        card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        card.line.width = Pt(1)
        # Top accent line
        _add_rect(s, x, y, Inches(2.9), Inches(0.08), color)
        # Emoji
        _add_textbox(
            s,
            x,
            y + Inches(0.4),
            Inches(2.9),
            Inches(0.7),
            emoji,
            font_size=36,
            alignment=PP_ALIGN.CENTER,
        )
        # Title
        _add_textbox(
            s,
            x + Inches(0.15),
            y + Inches(1.2),
            Inches(2.6),
            Inches(0.5),
            title,
            font_size=20,
            bold=True,
            color=color,
            alignment=PP_ALIGN.CENTER,
        )
        # Description
        _add_textbox(
            s,
            x + Inches(0.15),
            y + Inches(1.8),
            Inches(2.6),
            Inches(2.0),
            desc,
            font_size=15,
            color=DARK_TEXT,
            alignment=PP_ALIGN.CENTER,
        )

    _page_number(s, 11)

    # =================================================================
    # SLIDE 12 — Section: Como usar
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, PRIMARY_DARK)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(1),
        Inches(2.8),
        Inches(11),
        Inches(0.6),
        "4",
        font_size=18,
        color=PRIMARY_ACCENT,
        alignment=PP_ALIGN.CENTER,
    )
    _add_textbox(
        s,
        Inches(1),
        Inches(3.3),
        Inches(11),
        Inches(1),
        "Como usar — demos ao vivo",
        font_size=36,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
    )
    _page_number(s, 12)

    # =================================================================
    # SLIDE 13 — Demo 1: Tracing em 1 linha
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, WHITE)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.3),
        Inches(11),
        Inches(0.5),
        "Demo 1 — Tracing em 1 linha",
        font_size=28,
        bold=True,
        color=PRIMARY_DARK,
    )
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.85),
        Inches(11),
        Inches(0.4),
        "make tracing  ·  tracing_basics.py",
        font_size=16,
        color=PRIMARY_ACCENT,
    )

    _add_code_block(
        s,
        Inches(0.8),
        Inches(1.5),
        Inches(11.7),
        Inches(1.5),
        [
            "import mlflow",
            "",
            "mlflow.openai.autolog()   # <- 1 linha instrumenta todas as chamadas",
            "",
            "client = get_client()",
            "response = client.chat.completions.create(",
            "    model=MODEL_NAME,",
            "    messages=[...],",
            ")",
        ],
        font_size=13,
    )

    # Highlights
    tf = _add_rich_textbox(s, Inches(0.8), Inches(3.5), Inches(11.5), Inches(3.5))
    _add_bullet(
        tf,
        "👉 Captura inputs, outputs, tokens e latência automaticamente",
        size=18,
        space_before=Pt(8),
    )
    _add_bullet(tf, "👉 Sem nenhuma instrumentação manual", size=18, space_before=Pt(6))
    _add_bullet(
        tf,
        "👉 Abra o trace no MLflow UI e veja tudo pronto",
        size=18,
        space_before=Pt(6),
    )
    _add_bullet(
        tf,
        "⚠️ Gotcha: sempre chame mlflow.flush_trace_async_logging() "
        "antes de search_traces()",
        size=16,
        color=ACCENT_ORANGE,
        space_before=Pt(16),
    )

    _page_number(s, 13)

    # =================================================================
    # SLIDE 14 — Demo 2: Spans aninhados
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, WHITE)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.3),
        Inches(11),
        Inches(0.5),
        "Demo 2 — Pipelines complexos: spans aninhados",
        font_size=28,
        bold=True,
        color=PRIMARY_DARK,
    )
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.85),
        Inches(11),
        Inches(0.4),
        "make tracing  ·  tracing_basics.py",
        font_size=16,
        color=PRIMARY_ACCENT,
    )

    _add_code_block(
        s,
        Inches(0.8),
        Inches(1.5),
        Inches(11.7),
        Inches(2.5),
        [
            "@mlflow.trace",
            "def demo_rag_pipeline(question: str) -> str:",
            "    context = retrieve_context(question)        # span filho",
            "    answer = generate_answer(question, context) # span filho",
            "    return answer",
            "",
            '@mlflow.trace(span_type="RETRIEVER")',
            "def retrieve_context(question: str) -> str:",
            "    ...  # busca em vector store",
            "",
            '@mlflow.trace(span_type="LLM")',
            "def generate_answer(question: str, context: str) -> str:",
            "    ...  # chamada ao LLM com o contexto",
        ],
        font_size=13,
    )

    hbox = _add_rect(
        s, Inches(0.8), Inches(4.4), Inches(11.7), Inches(1.2), SECONDARY_ACCENT
    )
    tf = hbox.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(20)
    tf.margin_right = Pt(20)
    tf.margin_top = Pt(14)
    p = tf.paragraphs[0]
    _add_run(p, "💡 ", size=18)
    _add_run(
        p,
        "Quando a resposta final está errada, o trace mostra: o problema foi "
        "o retrieval (contexto ruim) ou a geração (modelo)?",
        size=16,
        color=DARK_TEXT,
    )

    _page_number(s, 14)

    # =================================================================
    # SLIDE 15 — Demo 3: Tokens + Cost
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, WHITE)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.3),
        Inches(11),
        Inches(0.5),
        "Demo 3 — Quanto isso custa?",
        font_size=28,
        bold=True,
        color=PRIMARY_DARK,
    )
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.85),
        Inches(11),
        Inches(0.4),
        "make tokens  ·  token_usage.py",
        font_size=16,
        color=PRIMARY_ACCENT,
    )

    # Two boxes side by side
    # Left: autolog info
    _add_rect(s, Inches(0.8), Inches(1.5), Inches(5.7), Inches(2.5), SECONDARY_ACCENT)
    tf = _add_rich_textbox(s, Inches(1.0), Inches(1.6), Inches(5.3), Inches(2.3))
    _add_bullet(
        tf,
        "✅ MLflow calcula custo automático",
        size=17,
        bold=True,
        color=PRIMARY_DARK,
        space_before=Pt(0),
    )
    _add_bullet(
        tf,
        "para OpenAI, Anthropic e provedores",
        size=17,
        bold=True,
        color=PRIMARY_DARK,
        space_before=Pt(0),
    )
    _add_bullet(
        tf,
        "com pricing registrado.",
        size=17,
        bold=True,
        color=PRIMARY_DARK,
        space_before=Pt(0),
    )

    # Right: manual
    _add_rect(
        s,
        Inches(6.8),
        Inches(1.5),
        Inches(5.7),
        Inches(2.5),
        RGBColor(0xFD, 0xF2, 0xE9),
    )
    tf2 = _add_rich_textbox(s, Inches(7.0), Inches(1.6), Inches(5.3), Inches(2.3))
    _add_bullet(
        tf2,
        "⚠️ Para modelos self-hosted:",
        size=17,
        bold=True,
        color=ACCENT_ORANGE,
        space_before=Pt(0),
    )
    _add_bullet(
        tf2,
        "atribua custo manualmente no span",
        size=17,
        bold=True,
        color=ACCENT_ORANGE,
        space_before=Pt(0),
    )

    _add_code_block(
        s,
        Inches(0.8),
        Inches(4.3),
        Inches(11.7),
        Inches(2.0),
        [
            'span.set_attribute("mlflow.llm.cost", {',
            '    "input_cost": input_tokens * CUSTOM_INPUT_COST_PER_TOKEN,',
            '    "output_cost": output_tokens * CUSTOM_OUTPUT_COST_PER_TOKEN,',
            '    "total_cost": ...',
            "})",
        ],
        font_size=14,
    )

    _add_textbox(
        s,
        Inches(0.8),
        Inches(6.5),
        Inches(11),
        Inches(0.4),
        "👉 Veja 'Cost Breakdown' e 'Token Usage' por trace no MLflow UI",
        font_size=16,
        bold=True,
        color=PRIMARY_ACCENT,
    )

    _page_number(s, 15)

    # =================================================================
    # SLIDE 16 — Demo 4: Judges
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, WHITE)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.3),
        Inches(11),
        Inches(0.5),
        "Demo 4 — A resposta estava certa?",
        font_size=28,
        bold=True,
        color=PRIMARY_DARK,
    )
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.85),
        Inches(11),
        Inches(0.4),
        "make judges  ·  judges.py",
        font_size=16,
        color=PRIMARY_ACCENT,
    )

    # LLM Judge card
    _add_rect(s, Inches(0.8), Inches(1.5), Inches(5.7), Inches(3.5), WHITE, True)
    _add_rect(s, Inches(0.8), Inches(1.5), Inches(5.7), Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(1.0),
        Inches(1.7),
        Inches(5.3),
        Inches(0.5),
        "🧠 LLM Judge",
        font_size=20,
        bold=True,
        color=PRIMARY_ACCENT,
    )
    _add_textbox(
        s,
        Inches(1.0),
        Inches(2.2),
        Inches(5.3),
        Inches(0.4),
        "Um modelo julga a resposta segundo regras\nem linguagem natural",
        font_size=14,
        color=RGBColor(0x66, 0x66, 0x66),
    )
    _add_code_block(
        s,
        Inches(1.0),
        Inches(2.8),
        Inches(5.3),
        Inches(1.8),
        [
            "Guidelines(",
            '    name="technical_accuracy",',
            "    guidelines=(",
            '        "A resposta deve ser tecnicamente',
            '         precisa sobre MLflow."',
            "    ),",
            "    model=JUDGE_MODEL,",
            ")",
        ],
        font_size=11,
    )

    # Code-based card
    _add_rect(s, Inches(6.8), Inches(1.5), Inches(5.7), Inches(3.5), WHITE, True)
    _add_rect(
        s,
        Inches(6.8),
        Inches(1.5),
        Inches(5.7),
        Inches(0.08),
        RGBColor(0x27, 0xAE, 0x60),
    )
    _add_textbox(
        s,
        Inches(7.0),
        Inches(1.7),
        Inches(5.3),
        Inches(0.5),
        "🔢 Code-based Scorer",
        font_size=20,
        bold=True,
        color=RGBColor(0x27, 0xAE, 0x60),
    )
    _add_textbox(
        s,
        Inches(7.0),
        Inches(2.2),
        Inches(5.3),
        Inches(0.4),
        "Regra determinística em Python — sem custo\nde LLM e sem latência",
        font_size=14,
        color=RGBColor(0x66, 0x66, 0x66),
    )
    _add_code_block(
        s,
        Inches(7.0),
        Inches(2.8),
        Inches(5.3),
        Inches(1.8),
        [
            "@scorer",
            "def no_hallucination_keywords(",
            "    inputs, outputs",
            ") -> Feedback:",
            '    red_flags = ["estudos mostram",',
            '                 "comprovado cientificamente"]',
            "    if any(f in str(outputs) for f in red_flags):",
            "        return Feedback(value=False)",
            "    return Feedback(value=True)",
        ],
        font_size=11,
    )

    # Tip
    hbox = _add_rect(
        s, Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.7), PRIMARY_DARK
    )
    tf = hbox.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(20)
    tf.margin_top = Pt(12)
    p = tf.paragraphs[0]
    _add_run(p, "💡 Dica: combine os dois. ", size=16, bold=True, color=WHITE)
    _add_run(
        p,
        "LLM judges medem qualidade subjetiva; scorers de código "
        "garantem regras objetivas de graça.",
        size=16,
        color=WHITE,
    )

    # Evaluate call
    _add_code_block(
        s,
        Inches(0.8),
        Inches(6.2),
        Inches(11.7),
        Inches(0.9),
        [
            "results = mlflow.genai.evaluate(data=dataset, predict_fn=predict_fn, scorers=[...])",
        ],
        font_size=12,
    )

    _page_number(s, 16)

    # =================================================================
    # SLIDE 17 — Demo 5: Streaming + span manual (conceitos)
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, WHITE)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.3),
        Inches(11),
        Inches(0.5),
        "Demo 5 — Agente real com streaming + span manual",
        font_size=28,
        bold=True,
        color=PRIMARY_DARK,
    )
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.85),
        Inches(11),
        Inches(0.4),
        "make langchain-agent  ·  langchain_agent.py",
        font_size=16,
        color=PRIMARY_ACCENT,
    )

    _add_textbox(
        s,
        Inches(0.8),
        Inches(1.4),
        Inches(11.5),
        Inches(0.5),
        "🏭 Mesmo padrão do backend de produção keepee-rag RAG",
        font_size=18,
        bold=True,
        color=PRIMARY_DARK,
    )

    tf = _add_rich_textbox(s, Inches(0.8), Inches(2.1), Inches(11.5), Inches(4.5))
    for item in [
        "agent.astream() — streaming token a token em tempo real",
        "Span manual SpanType.AGENT com inputs/outputs explícitos",
        "Session management via dict (reconstrução manual do histórico)",
        "get_stream_writer() para logs de progresso das tools",
        "Eventos JSON newline-delimited via make_event()",
        "trace_id capturado do span e exposto para anexar feedback",
        "User/session vinculados ao trace ANTES do stream",
        "Tags (provider, model_name) setadas após o stream",
    ]:
        _add_bullet(tf, f"► {item}", size=17, space_before=Pt(9))

    _page_number(s, 17)

    # =================================================================
    # SLIDE 18 — Demo 5: autolog vs manual (tabela)
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, WHITE)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.4),
        Inches(11),
        Inches(0.6),
        "autolog vs. Manual Span",
        font_size=28,
        bold=True,
        color=PRIMARY_DARK,
    )

    rows, cols = 7, 3
    tbl_shape = s.shapes.add_table(
        rows, cols, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.2)
    )
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(2.5)
    tbl.columns[1].width = Inches(4.6)
    tbl.columns[2].width = Inches(4.6)

    for ci, h in enumerate(
        ["Aspecto", "mlflow.langchain.autolog()", "Padrão produção (manual span)"]
    ):
        cell = tbl.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_DARK

    tdata = [
        ["Streaming", "invoke() bloqueante", "astream() token a token"],
        ["Session", "MemorySaver (checkpointer)", "Dict (reconstrução manual)"],
        ["Controle do span", "Automático (genérico)", "Inputs/outputs explícitos"],
        ["Progresso tools", "Só o resultado final", "Logs via get_stream_writer()"],
        ["Eventos", "Apenas no trace", "JSON newline-delimited (SSE)"],
        ["trace_id", "get_last_active_trace_id()", "Capturado do span manual"],
    ]
    for ri, row in enumerate(tdata):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.color.rgb = DARK_TEXT
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER if ci == 0 else PP_ALIGN.LEFT
                p.space_before = Pt(4)
                p.space_after = Pt(4)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            if ci == 0:
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True

    _page_number(s, 18)

    # =================================================================
    # SLIDE 19 — Demo 6: Feedback + Sampling
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, WHITE)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.3),
        Inches(11),
        Inches(0.5),
        "Demo 6 — Operando em produção",
        font_size=28,
        bold=True,
        color=PRIMARY_DARK,
    )
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.85),
        Inches(11),
        Inches(0.4),
        "make monitoring  ·  production_monitoring.py",
        font_size=16,
        color=PRIMARY_ACCENT,
    )

    # Sampling section
    _add_textbox(
        s,
        Inches(0.8),
        Inches(1.4),
        Inches(5),
        Inches(0.4),
        "📊 Sampling diferenciado por criticidade",
        font_size=20,
        bold=True,
        color=PRIMARY_ACCENT,
    )
    _add_textbox(
        s,
        Inches(0.8),
        Inches(1.85),
        Inches(11.5),
        Inches(0.4),
        "Não traceie 100% — storage e custo importam. Use sampling_ratio_override:",
        font_size=15,
        color=DARK_TEXT,
    )
    _add_code_block(
        s,
        Inches(0.8),
        Inches(2.3),
        Inches(11.7),
        Inches(1.4),
        [
            "@mlflow.trace(sampling_ratio_override=1.0)  # 100% — pagamentos (crítico)",
            "@mlflow.trace(sampling_ratio_override=0.1)  # 10%  — alto volume",
            "def agent_call(agent, query, user_id, session_id):",
            "    return agent_invoke(agent, query, user_id, session_id)",
        ],
        font_size=13,
    )

    # Feedback section
    _add_textbox(
        s,
        Inches(0.8),
        Inches(4.0),
        Inches(5),
        Inches(0.4),
        "👍👎 Feedback humano no trace",
        font_size=20,
        bold=True,
        color=PRIMARY_ACCENT,
    )
    _add_code_block(
        s,
        Inches(0.8),
        Inches(4.5),
        Inches(11.7),
        Inches(1.7),
        [
            "trace_id = mlflow.get_last_active_trace_id()",
            "mlflow.log_feedback(",
            "    trace_id=trace_id,",
            '    name="user_rating",',
            "    value=True,       # 👍 ou 👎",
            '    rationale="texto livre do usuário",',
            '    source=AssessmentSource(source_type="HUMAN"),',
            ")",
        ],
        font_size=12,
    )

    hbox = _add_rect(
        s,
        Inches(0.8),
        Inches(6.5),
        Inches(11.7),
        Inches(0.5),
        RGBColor(0xFD, 0xF2, 0xE9),
    )
    tf = hbox.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(16)
    tf.margin_top = Pt(6)
    p = tf.paragraphs[0]
    _add_run(
        p,
        "⚠️  flush_trace_async_logging() antes de search_traces()!",
        size=15,
        bold=True,
        color=ACCENT_ORANGE,
    )

    _page_number(s, 19)

    # =================================================================
    # SLIDE 20 — Section: Fechamento
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, PRIMARY_DARK)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(1),
        Inches(2.8),
        Inches(11),
        Inches(0.6),
        "5",
        font_size=18,
        color=PRIMARY_ACCENT,
        alignment=PP_ALIGN.CENTER,
    )
    _add_textbox(
        s,
        Inches(1),
        Inches(3.3),
        Inches(11),
        Inches(1),
        "Fechamento",
        font_size=36,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
    )
    _page_number(s, 20)

    # =================================================================
    # SLIDE 21 — Ciclo virtuoso
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, WHITE)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.4),
        Inches(11),
        Inches(0.6),
        "O ciclo virtuoso da observabilidade",
        font_size=28,
        bold=True,
        color=PRIMARY_DARK,
    )

    # Flow boxes in a circle
    flow = [
        ("🔍 Tracing", "Vejo o que\nacontece", PRIMARY_ACCENT),
        ("✅ Avaliação", "Meço a\nqualidade", RGBColor(0x27, 0xAE, 0x60)),
        ("⚡ Otimização", "Melhoro prompts\ne tools", ACCENT_ORANGE),
    ]

    starts = [Inches(1.5), Inches(5.5), Inches(9.5)]
    for i, (emoji_title, desc, color) in enumerate(flow):
        x = starts[i]
        y = Inches(1.8)
        card = _add_rect(s, x, y, Inches(3.2), Inches(2.8), WHITE, line=True)
        card.line.color.rgb = color
        card.line.width = Pt(2)
        _add_rect(s, x, y, Inches(3.2), Inches(0.08), color)
        _add_textbox(
            s,
            x + Inches(0.1),
            y + Inches(0.3),
            Inches(3.0),
            Inches(0.5),
            emoji_title,
            font_size=22,
            bold=True,
            color=color,
            alignment=PP_ALIGN.CENTER,
        )
        _add_textbox(
            s,
            x + Inches(0.1),
            y + Inches(1.0),
            Inches(3.0),
            Inches(1.2),
            desc,
            font_size=17,
            color=DARK_TEXT,
            alignment=PP_ALIGN.CENTER,
        )

    # Arrows between boxes (simple text arrows)
    _add_textbox(
        s,
        Inches(4.1),
        Inches(2.7),
        Inches(2),
        Inches(0.5),
        "➡️",
        font_size=36,
        color=PRIMARY_ACCENT,
        alignment=PP_ALIGN.CENTER,
    )
    _add_textbox(
        s,
        Inches(8.1),
        Inches(2.7),
        Inches(2),
        Inches(0.5),
        "➡️",
        font_size=36,
        color=PRIMARY_ACCENT,
        alignment=PP_ALIGN.CENTER,
    )

    # Curved arrow back
    _add_textbox(
        s,
        Inches(5.5),
        Inches(4.8),
        Inches(3),
        Inches(0.5),
        "⬆️  fecha o ciclo",
        font_size=16,
        bold=True,
        color=PRIMARY_ACCENT,
        alignment=PP_ALIGN.CENTER,
    )

    # Bottom message
    _add_rect(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.7), PRIMARY_DARK)
    _add_textbox(
        s,
        Inches(1.0),
        Inches(5.9),
        Inches(11.3),
        Inches(0.6),
        'Observabilidade não é só "ver logs" — é o que fecha o ciclo '
        "entre observar, medir e melhorar sua IA.",
        font_size=16,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
    )

    _page_number(s, 21)

    # =================================================================
    # SLIDE 22 — Checklist
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, WHITE)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.4),
        Inches(11),
        Inches(0.6),
        "Checklist para levar para casa",
        font_size=28,
        bold=True,
        color=PRIMARY_DARK,
    )

    checks = [
        "Ligue auto-tracing (mlflow.openai.autolog / mlflow.langchain.autolog)",
        "Use spans aninhados em pipelines (RAG, agentes) para isolar a etapa que falha",
        "Atribua custo manualmente se o modelo for self-hosted",
        "Combine LLM judges (qualidade) + code-based scorers (regras objetivas)",
        "Vincule user_id e session_id aos traces",
        "Em produção: sampling por criticidade + coleta de feedback humano",
        "Lembre do flush_trace_async_logging() antes de ler traces",
    ]

    tf = _add_rich_textbox(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5))
    for item in checks:
        _add_bullet(tf, f"☐  {item}", size=18, space_before=Pt(12))

    _page_number(s, 22)

    # =================================================================
    # SLIDE 23 — Comandos
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, WHITE)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(0.8),
        Inches(0.4),
        Inches(11),
        Inches(0.6),
        "Comandos para explorar",
        font_size=28,
        bold=True,
        color=PRIMARY_DARK,
    )

    _add_code_block(
        s,
        Inches(0.8),
        Inches(1.3),
        Inches(5.5),
        Inches(5.2),
        [
            "# Setup",
            "uv sync",
            "",
            "# Demos (make ou uv run)",
            "make tracing            # 01",
            "make tokens             # 02",
            "make judges             # 05",
            "make langchain-agent    # 11",
            "make monitoring         # 07",
            "",
            "make help    # lista todas",
        ],
        font_size=14,
    )

    # Right side: MLflow info
    _add_rect(s, Inches(6.8), Inches(1.3), Inches(5.7), Inches(2.5), SECONDARY_ACCENT)
    tf = _add_rich_textbox(s, Inches(7.0), Inches(1.5), Inches(5.3), Inches(2.2))
    _add_bullet(
        tf, "MLflow UI", size=20, bold=True, color=PRIMARY_DARK, space_before=Pt(0)
    )
    _add_bullet(
        tf, "http://localhost:5000", size=16, color=PRIMARY_ACCENT, space_before=Pt(4)
    )
    _add_bullet(tf, "", size=8, space_before=Pt(4))
    _add_bullet(
        tf,
        "MLflow AI Gateway",
        size=20,
        bold=True,
        color=PRIMARY_DARK,
        space_before=Pt(8),
    )
    _add_bullet(
        tf,
        "http://localhost:5000/gateway/mlflow/v1",
        size=16,
        color=PRIMARY_ACCENT,
        space_before=Pt(4),
    )
    _add_bullet(tf, "", size=8, space_before=Pt(4))
    _add_bullet(
        tf,
        "Docker Compose:",
        size=20,
        bold=True,
        color=PRIMARY_DARK,
        space_before=Pt(8),
    )
    _add_bullet(
        tf, "podman compose up -d", size=16, color=PRIMARY_ACCENT, space_before=Pt(4)
    )

    # Make info
    _add_textbox(
        s,
        Inches(6.8),
        Inches(4.2),
        Inches(5.7),
        Inches(0.4),
        "Cada demo cria um experiment numerado no MLflow:",
        font_size=15,
        color=DARK_TEXT,
    )
    _add_code_block(
        s,
        Inches(6.8),
        Inches(4.7),
        Inches(5.7),
        Inches(1.8),
        [
            "01 - tracing_basics",
            "02 - token_usage",
            "05 - judges",
            "11 - langchain_agent",
            "07 - production_monitoring",
        ],
        font_size=13,
    )

    _page_number(s, 23)

    # =================================================================
    # SLIDE 24 — Obrigado / Referências
    # =================================================================
    s = new_slide()
    _set_slide_bg(s, PRIMARY_DARK)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY_ACCENT)
    _add_textbox(
        s,
        Inches(1),
        Inches(1.5),
        Inches(11),
        Inches(1),
        "Obrigado!",
        font_size=44,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
    )
    _add_rect(s, Inches(5), Inches(2.6), Inches(3.333), Inches(0.03), PRIMARY_ACCENT)

    _add_textbox(
        s,
        Inches(1),
        Inches(3.0),
        Inches(11),
        Inches(0.5),
        "Referências",
        font_size=20,
        bold=True,
        color=PRIMARY_ACCENT,
        alignment=PP_ALIGN.CENTER,
    )

    refs = [
        "MLflow GenAI Docs — mlflow.org/docs/latest/genai",
        "Tracing Quickstart — mlflow.org/docs/latest/genai/tracing/quickstart",
        "Evaluation & Monitoring — mlflow.org/docs/latest/genai/eval-monitor",
        "LLM Judges / Scorers — mlflow.org/docs/latest/genai/eval-monitor/scorers",
        "Production Monitoring — mlflow.org/docs/latest/genai/tracing/prod-tracing",
    ]
    tf = _add_rich_textbox(s, Inches(2), Inches(3.7), Inches(9.5), Inches(3.0))
    for ref in refs:
        _add_bullet(tf, f"📎  {ref}", size=16, color=WHITE, space_before=Pt(10))

    _page_number(s, 24)

    # -----------------------------------------------------------------
    # Salvar
    # -----------------------------------------------------------------
    output = "docs/workshop-observabilidade-ia.pptx"
    prs.save(output)
    print(f"✅ Apresentacao salva em: {output}")


if __name__ == "__main__":
    build()
